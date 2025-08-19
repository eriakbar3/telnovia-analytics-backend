# In telnovia-analytics-backend/app/presentation/router.py

from fastapi import APIRouter, Depends, HTTPException
from fastapi.responses import FileResponse
from sqlalchemy.orm import Session
import os
from pptx import Presentation
from pydantic import BaseModel
from typing import List
import uuid
from pptx.util import Inches

from app.database import get_db
from app.auth.oauth2 import get_current_user
from app import crud, schemas

router = APIRouter(
    prefix="/api/v1/notebook",
    tags=['Presentation']
)

@router.get("/{notebook_id}/presentation/preview", response_model=schemas.PresentationPreviewResponse)
def get_presentation_preview(
    notebook_id: int,
    db: Session = Depends(get_db),
    current_user: schemas.UserOut = Depends(get_current_user)
):
    """
    Menghasilkan representasi JSON dari slide presentasi untuk diedit di frontend.
    """
    notebook = crud.get_notebook(db, notebook_id=notebook_id, owner_id=current_user.id)
    if not notebook:
        raise HTTPException(status_code=404, detail="Notebook tidak ditemukan.")

    slides = []

    # Slide 1: Judul
    slides.append({
        "id": f"slide-{uuid.uuid4()}",
        "type": "title",
        "title": f"Analisis Data: {notebook.filename}",
        "content": f"Dibuat oleh Telnovia untuk {current_user.email}"
    })

    # Slide 2: Data Health Report
    if notebook.health_report and notebook.health_report.get("columns"):
        slides.append({
            "id": f"slide-{uuid.uuid4()}",
            "type": "table",
            "title": "Data Health Report",
            "content": notebook.health_report["columns"]
        })

    # Slide 3..N: Hasil Analisis dari Chat
    for conv in notebook.conversations:
        # Cek sederhana jika respons mengandung format tabel markdown
        if '|' in conv.ai_response and '---' in conv.ai_response:
            slides.append({
                "id": f"slide-{uuid.uuid4()}",
                "type": "markdown",
                "title": conv.user_query,
                "content": conv.ai_response
            })

    return {"slides": slides}

# Definisikan skema untuk request body yang baru
class FinalizePresentationRequest(BaseModel):
    slides: List[schemas.SlideBase]

@router.post("/{notebook_id}/presentation")
def create_presentation_for_notebook(
    notebook_id: int,
    request_body: FinalizePresentationRequest, # <-- Menerima body request
    db: Session = Depends(get_db),
    current_user: schemas.UserOut = Depends(get_current_user)
):
    """
    Membuat file presentasi .pptx berdasarkan data slide yang dikirim dari frontend.
    """
    notebook = crud.get_notebook(db, notebook_id=notebook_id, owner_id=current_user.id)
    if not notebook:
        raise HTTPException(status_code=404, detail="Notebook tidak ditemukan.")

    prs = Presentation()
    slides_data = request_body.slides

    # Membuat presentasi berdasarkan data JSON yang diterima dari frontend
    for slide_data in slides_data:
        if slide_data.type == "title":
            slide = prs.slides.add_slide(prs.slide_layouts[0])
            slide.shapes.title.text = slide_data.title
            if isinstance(slide_data.content, str):
                slide.placeholders[1].text = slide_data.content
        elif slide_data.type == "table" and isinstance(slide_data.content, list):
            slide = prs.slides.add_slide(prs.slide_layouts[5])
            slide.shapes.title.text = slide_data.title
            
            report_data = slide_data.content
            rows, cols = len(report_data) + 1, 4 # +1 untuk header
            table_shape = slide.shapes.add_table(rows, cols, Inches(0.5), Inches(1.5), Inches(9), Inches(5))
            table = table_shape.table

            # Header (diasumsikan struktur health report)
            table.cell(0, 0).text = 'Nama Kolom'
            table.cell(0, 1).text = 'Tipe Data'
            table.cell(0, 2).text = 'Nilai Kosong (%)'
            table.cell(0, 3).text = 'Peringatan'

            # Isi tabel
            for i, row_data in enumerate(report_data):
                table.cell(i + 1, 0).text = row_data.get('name', '')
                table.cell(i + 1, 1).text = row_data.get('dtype', '')
                table.cell(i + 1, 2).text = str(row_data.get('missing_values', {}).get('percentage', '0'))
                table.cell(i + 1, 3).text = ', '.join(row_data.get('warnings', []))
        elif slide_data.type == "markdown":
            slide = prs.slides.add_slide(prs.slide_layouts[5])
            slide.shapes.title.text = slide_data.title
            
            textbox = slide.shapes.add_textbox(Inches(1), Inches(1.5), Inches(8), Inches(4))
            tf = textbox.text_frame
            tf.text = str(slide_data.content) # Menampilkan markdown mentah

    # Simpan file presentasi dan kembalikan
    output_dir = "presentations"
    os.makedirs(output_dir, exist_ok=True)
    output_path = os.path.join(output_dir, f"{notebook_id}_final_presentation.pptx")
    prs.save(output_path)
    
    return FileResponse(path=output_path, media_type='application/vnd.openxmlformats-officedocument.presentationml.presentation', filename=f"Telnovia_Analysis_{notebook.filename}.pptx")